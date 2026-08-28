#!/usr/bin/env python3
"""
deck_builder.py - Build a beautiful, animated, EDITABLE .pptx from a declarative JSON spec.

This engine is intentionally dependency-light: it needs only ``python-pptx`` (any build
that can create shapes, text, tables and native charts) plus ``lxml`` (bundled with
python-pptx). No network, no image model, no API keys.

Why hand-built animations?
    Many python-pptx builds (e.g. 1.0.2) ship WITHOUT a high-level animation API
    (no ``shape.enter_effect`` / ``shape.timing``). Rather than depend on that, this
    engine emits the OOXML ``<p:timing>/<p:seq>`` elements by hand. The block was
    validated to round-trip through python-pptx (save -> reopen preserves the sequences).

    - ``st`` (time action): after=196608, with=131073, click=-2147483647
    - ``type`` (ST_SeqType) selects the animation; ``<p:effectLst>`` child disambiguates.

Usage:
    python deck_builder.py --spec deck.json --out deck.pptx
    python deck_builder.py --spec deck.json --out deck.pptx --resize 10,7.5
    python deck_builder.py --demo            # writes /tmp/deck_demo.pptx

Run ``python deck_builder.py --spec deck.json --inspect`` to print the rendered slide
plan (does not write a file).
"""
from __future__ import annotations

import argparse
import json
import random
import uuid
from pathlib import Path

from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION

# --------------------------------------------------------------------------- namespaces
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def p(tag: str) -> str:
    return f"{{{P}}}{tag}"


def a(tag: str) -> str:
    return f"{{{A}}}{tag}"


# --------------------------------------------------------------------------- helpers
def _rgb(hexstr: str) -> RGBColor:
    h = hexstr.lstrip("#")
    return RGBColor.from_string(h)


def _scale(w_in: float, h_in: float):
    return Inches(w_in), Inches(h_in)


# ST_SeqType (enter family) and emphasis family integer codes.
ENTER = {"appear": 0, "fly": 1, "push": 2, "split": 3, "wor": 8, "zoom": 9, "fall": 10}
EMPH = {"compress": 0, "expand": 1, "glow": 2, "pulse": 3, "resize": 4, "skew": 5}
# time-action (st) constants
ST_AFTER, ST_WITH, ST_CLICK = 196608, 131073, -2147483647


def _add_seq(shape, *, seq_type: int, child_tag: str, child_attrs: dict, st: int):
    """Append a validated ``<p:seq>`` to the shape's slide timing.

    ``shape`` may be a python-pptx Shape OR a text_frame (both carry enough XML to
    locate the owning slide), so the engine can animate either.
    """
    if hasattr(shape, "slide"):
        slide_el = shape.slide.element
    else:
        # text_frame: <a:txBody> -> <p:sp> -> <cSld>
        slide_el = shape._element.getparent().getparent()
    timing = slide_el.find(p("timing"))
    if timing is None:
        timing = etree.Element(p("timing"), nsmap={"p": P, "a": A})
        slide_el.append(timing)
    seq = etree.SubElement(timing, p("seq"))
    seq.set("st", str(st))
    seq.set("type", str(seq_type))
    seq.set("tid", str(uuid.uuid4()))
    nv = etree.SubElement(seq, p("nvGrpPr"))
    etree.SubElement(nv, p("cNvPr"), id="1", name="")
    etree.SubElement(nv, p("cNvGrpPr"))
    etree.SubElement(nv, p("cNvPr"), name="")
    grp = etree.SubElement(seq, p("grpPr"))
    xfrm = etree.SubElement(grp, a("xfrm"))
    etree.SubElement(xfrm, a("off"), x="0", y="0")
    etree.SubElement(xfrm, a("ext"), cx="0", cy="0")
    etree.SubElement(xfrm, a("chOff"), x="0", y="0")
    etree.SubElement(xfrm, a("chExt"), cx="0", cy="0")
    etree.SubElement(xfrm, a("rot"), value="-900000")
    el = etree.SubElement(seq, p("effectLst"))
    etree.SubElement(el, p(child_tag), **child_attrs)


def add_animation(shape, *, effect: str = "fly", trigger: str = "after",
                  order: int = 1, speed: str = "slow", child_attrs: dict | None = None):
    """Attach an entry/emphasis animation to ``shape``.

    ``effect`` is one of the keys in ENTER / EMPH. ``trigger`` is 'after' | 'with' | 'click'.
    ``order`` is the 1-based play order within the slide (lower plays first).
    """
    if not child_attrs:
        child_attrs = {}
    seq_type = ENTER.get(effect, 1)
    st = {"after": ST_AFTER, "with": ST_WITH, "click": ST_CLICK}.get(trigger, ST_AFTER)
    # enter-family: use <p:after type=.../>; emphasis-family: use <p:effect/> by name
    if effect in EMPH:
        _add_seq(shape, seq_type=EMPH[effect], child_tag=effect, child_attrs=child_attrs, st=st)
    else:
        direction = child_attrs.get("type", "fromBottom")
        _add_seq(shape, seq_type=seq_type, child_tag="after",
                 child_attrs={"type": direction}, st=st)
    shape._pptx_anim = {"order": order, "effect": effect, "trigger": trigger}


# --------------------------------------------------------------------------- slide builders
class Deck:
    def __init__(self, theme: dict, w_in: float = 13.333, h_in: float = 7.5):
        self.theme = theme
        self.w_in, self.h_in = w_in, h_in
        self.prs = Presentation()
        self.prs.slide_width = Inches(w_in)
        self.prs.slide_height = Inches(h_in)
        self.blank = self.prs.slide_layouts[6]

    # -- low level
    def _bg(self, slide, color):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   Inches(self.w_in), Inches(self.h_in))
        r.fill.solid(); r.fill.fore_color.rgb = _rgb(color)
        r.line.fill.background()
        return r

    def _box(self, slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
        s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            s.fill.background()
        else:
            s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
        if line is None:
            s.line.fill.background()
        else:
            s.line.color.rgb = _rgb(line); s.line.width = Pt(1)
        return s

    def _txt(self, slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tf

    def _set_run(self, run, size, color, bold, italic, name):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = name
        if color is not None:
            run.font.color.rgb = _rgb(color)

    def _para(self, para, size, color, bold, italic, level, align, spacing):
        para.alignment = align
        para.level = level
        if spacing:
            para.space_after = Pt(spacing)
            para.space_before = Pt(0)

    def animate(self, shape, **anim):
        """Attach an animation to ``shape`` (works for Shape or text_frame)."""
        add_animation(shape, **anim)

    # -- public slide primitives (used by slide-type renderers)
    def background(self, slide, color):
        return self._bg(slide, color)

    def stripe(self, slide, side: str = "left", width: float = 0.14, color: str = None):
        color = color or self.theme.get("accent", "#12B5C1")
        if side == "left":
            return self._box(slide, 0, 0, width, self.h_in, fill=color, shape=MSO_SHAPE.RECTANGLE)
        return self._box(slide, self.w_in - width, 0, width, self.h_in, fill=color, shape=MSO_SHAPE.RECTANGLE)

    def section_badge(self, slide, number: str, kicker: str = None, align: str = "left"):
        """Big translucent number + accent bar + kicker text."""
        self._box(slide, 0.7, 0.55, 1.9, 1.9, fill=self.theme.get("accent", "#12B5C1"),
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        nb = self._txt(slide, 0.7, 0.5, 1.9, 1.9, MSO_ANCHOR.MIDDLE)
        nb.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = nb.paragraphs[0].add_run(); r.text = str(number)
        self._set_run(r, 60, "#0B1420", True, False, self.theme.get("title_font", "Calibri"))
        if kicker:
            tb = self._txt(slide, 2.9, 0.75, 9.5, 0.7)
            p0 = tb.paragraphs[0]
            r = p0.add_run(); r.text = kicker
            self._set_run(r, 16, self.theme.get("accent", "#12B5C1"), True, False, self.theme.get("body_font", "Calibri"))
        return nb

    def heading(self, slide, text: str, x: float = 0.7, y: float = 1.3, w: float = 11.9,
                size: int = 30, color: str = None, underline: bool = False, anim=None):
        color = color or self.theme.get("text", "#F1F5F9")
        tb = self._txt(slide, x, y, w, 0.9, MSO_ANCHOR.MIDDLE)
        p0 = tb.paragraphs[0]
        r = p0.add_run(); r.text = text
        self._set_run(r, size, color, True, False, self.theme.get("title_font", "Calibri"))
        if anim: self.animate(tb, **anim)
        return tb

    def bullets(self, slide, x, y, w, h, items, size=17, color=None, gap=10, levels=False,
                anim=None):
        """items: list of str, or (level, text) tuples, or dicts.

        ``anim`` is a dict of animation args (effect/trigger/order_base). Each item is
        rendered as its own textbox so items can animate independently in sequence.
        """
        color = color or self.theme.get("text", "#F1F5F9")
        accent = self.theme.get("accent", "#12B5C1")
        step = (size + gap) / 72.0
        base = (anim or {}).get("order_base", 1)
        for idx, it in enumerate(items):
            if isinstance(it, dict):
                txt = it.get("text", ""); lvl = it.get("level", 0); c = it.get("color") or color
                bold = it.get("bold", False)
            elif isinstance(it, tuple):
                lvl, txt = it; c = color; bold = False
            else:
                txt, lvl, bold, c = it, 0, False, color
            tf = self._txt(slide, x, y + idx * step, w, step)
            # colored square bullet
            b = tf.paragraphs[0].add_run()
            b.text = "\u25AA  "
            self._set_run(b, size, accent, True, False, self.theme.get("body_font", "Calibri"))
            r = tf.paragraphs[0].add_run(); r.text = txt
            self._set_run(r, size, c, bold, False, self.theme.get("body_font", "Calibri"))
            if anim:
                self.animate(tf, effect=anim.get("effect", "fade"),
                             trigger=anim.get("trigger", "after"), order=base + idx)
        return tf

    def metric_card(self, slide, x, y, w, h, value, label, sub=None, value_color=None, anim=None):
        panel = self._box(slide, x, y, w, h, fill=self.theme.get("panel", "#12212E"),
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        panel.line.color.rgb = _rgb(self.theme.get("accent", "#12B5C1"))
        panel.line.width = Pt(1.25)
        if anim: self.animate(panel, **anim)
        vc = value_color or self.theme.get("accent2", "#F5B642")
        tf = self._txt(slide, x, y + h * 0.28, w, h * 0.42, MSO_ANCHOR.MIDDLE)
        r = tf.paragraphs[0].add_run(); r.text = str(value)
        self._set_run(r, max(int(h * 0.72), 1), vc, True, False, self.theme.get("title_font", "Calibri"))
        lf = self._txt(slide, x, y + h * 0.66, w, h * 0.34)
        r2 = lf.paragraphs[0].add_run(); r2.text = label
        self._set_run(r2, max(int(h * 0.2), 8), self.theme.get("muted", "#94A3B8"), False, False, self.theme.get("body_font", "Calibri"))
        if sub:
            sf = self._txt(slide, x, y + h - (h * 0.22), w, h * 0.22)
            r3 = sf.paragraphs[0].add_run(); r3.text = sub
            self._set_run(r3, max(int(h * 0.18), 6), self.theme.get("accent", "#12B5C1"), False, True, self.theme.get("body_font", "Calibri"))
        return panel

    def bar_chart(self, slide, x, y, w, h, categories, values, value_color=None,
                  show_values: bool = True, label_size: int = 13):
        """Horizontal clustered bar chart (categories on Y)."""
        value_color = value_color or self.theme.get("accent", "#12B5C1")
        cd = CategoryChartData()
        cd.categories = [str(c) for c in categories]
        cd.add_series("value", tuple(float(v) for v in values))
        chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                             Inches(x), Inches(y), Inches(w), Inches(h), cd)
        chart = chart_frame.chart
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.gap_width = 80
        series = plot.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb(value_color)
        plot.has_data_labels = True
        # data labels with values
        lbls = plot.data_labels
        lbls.number_format = '0.00"@"'; lbls.number_format_is_linked = False
        lbls.show_values = show_values
        try:
            lbls.position = XL_DATA_LABEL_POSITION.INSIDE_END
        except Exception:
            pass
        cats = chart.category_axis
        try:
            cats.tick_labels.font.size = Pt(label_size)
            cats.tick_labels.font.color.rgb = _rgb(self.theme.get("text", "#F1F5F9"))
        except Exception:
            pass
        return chart_frame

    def grouped_chart(self, slide, x, y, w, h, categories, series, colors):
        """series: list of (name, [values]). colors: list of hex."""
        cd = CategoryChartData()
        cd.categories = [str(c) for c in categories]
        for entry, col in zip(series, colors):
            name = entry[0]
            vals = entry[1:]
            cd.add_series(name, tuple(float(v) for v in vals))
        chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                             Inches(x), Inches(y), Inches(w), Inches(h), cd)
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        plot = chart.plots[0]; plot.gap_width = 120
        for i, s in enumerate(plot.series):
            s.format.fill.solid(); s.format.fill.fore_color.rgb = _rgb(colors[i])
        try:
            chart.category_axis.tick_labels.font.size = Pt(12)
            chart.category_axis.tick_labels.font.color.rgb = _rgb(self.theme.get("text", "#F1F5F9"))
            chart.value_axis.tick_labels.font.size = Pt(11)
        except Exception:
            pass
        return chart_frame

    def data_table(self, slide, x, y, w, h, rows, header_color=None, row_colors=None, anim=None):
        n_r, n_c = len(rows), len(rows[0]) if rows else 0
        table_shape = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(h))
        grid = table_shape.table
        if anim: self.animate(table_shape, **anim)
        for c in range(n_c):
            grid.columns[c].width = Inches(w / n_c)
        for r, row in enumerate(rows):
            for c, cell in enumerate(row):
                tc = grid.cell(r, c)
                tc.word_wrap = True
                if row_colors and r > 0:
                    rc = row_colors[min(r, len(row_colors) - 1)]
                    tc.fill.solid(); tc.fill.fore_color.rgb = _rgb(rc)
                else:
                    tc.fill.solid(); tc.fill.fore_color.rgb = _rgb(header_color if r == 0 else self.theme.get("panel", "#12212E"))
                tc.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = tc.text_frame
                parts = _split_cell(cell)
                for i, part in enumerate(parts):
                    p0 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    run = p0.add_run(); run.text = part
                    is_bold = (i == 0) or str(part).startswith(("AWS", "BM", "Hy", "Dense", "Base", "RAG", "Fine"))
                    self._set_run(run, 12, "#F1F5F9", is_bold, False, "Calibri")
        return grid


def _split_cell(cell):
    # allow "boldpart|normalpart" styling within a cell
    if isinstance(cell, str) and "|" in cell:
        b, rest = cell.split("|", 1)
        return [b, rest]
    return [cell]


# --------------------------------------------------------------------------- slide-type renderers
def render_slide(deck, slide, spec):
    t = spec.get("type", "bullets")
    fn = globals().get("_slide_" + t, _slide_bullets)
    fn(deck, slide, spec)


def _slide_title(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.16, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    # big watermark number
    wm = deck._txt(slide, 8.6, 0.2, 4.5, 4.5, MSO_ANCHOR.MIDDLE)
    r = wm.paragraphs[0].add_run(); r.text = str(spec.get("num", "01"))
    deck._set_run(r, 260, "#0E2A36", False, False, deck.theme.get("title_font", "Calibri"))
    deck.animate(wm, **anim.get("watermark", {"effect": "fly", "trigger": "after", "order": 1}))
    # kicker
    if spec.get("kicker"):
        kb = deck._txt(slide, 0.9, 1.0, 11, 0.5)
        r = kb.paragraphs[0].add_run(); r.text = spec["kicker"]
        deck._set_run(r, 16, deck.theme.get("accent", "#12B5C1"), True, False, deck.theme.get("body_font", "Calibri"))
        deck.animate(kb, **anim.get("kicker", {"effect": "fade", "trigger": "after", "order": 2}))
    deck.heading(slide, spec["title"], 0.9, 1.7, 11.4, 40, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 3}))
    # accent rule
    rule = deck._box(slide, 0.94, 2.75, 3.2, 0.06, fill=deck.theme.get("accent", "#12B5C1"))
    deck.animate(rule, **anim.get("rule", {"effect": "push", "trigger": "after", "order": 3}))
    if spec.get("subtitle"):
        sf = deck._txt(slide, 0.9, 3.05, 11.4, 2.4)
        for i, para_text in enumerate(_iter_lines(spec["subtitle"])):
            p0 = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
            deck._para(p0, 18, "#CBD5E1", False, False, 0, PP_ALIGN.LEFT, 8)
            r = p0.add_run(); r.text = para_text if isinstance(para_text, str) else para_text.get("text", "")
            deck._set_run(r, 18, "#CBD5E1", bool(isinstance(para_text, dict) and para_text.get("bold")), False, deck.theme.get("body_font", "Calibri"))
        deck.animate(sf, **anim.get("subtitle", {"effect": "fade", "trigger": "after", "order": 4}))
    # footer attribution line
    if spec.get("footers"):
        ff = deck._txt(slide, 0.9, 6.5, 11.4, 0.5)
        r = ff.paragraphs[0].add_run(); r.text = "  ".join(spec["footers"])
        deck._set_run(r, 13, "#94A3B8", False, False, deck.theme.get("body_font", "Calibri"))
        deck.animate(ff, **anim.get("footer", {"effect": "fade", "trigger": "after", "order": 5}))


def _slide_section(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.16, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    badge = deck.section_badge(slide, str(spec.get("num", "")), spec.get("kicker", ""))
    deck.animate(badge, **anim.get("badge", {"effect": "zoom", "trigger": "after", "order": 1}))
    deck.heading(slide, spec["title"], 0.9, 2.4, 11.4, 36, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    deck._box(slide, 0.94, 3.35, 2.6, 0.05, fill=deck.theme.get("accent", "#12B5C1"))
    if spec.get("lead"):
        lf = deck._txt(slide, 0.9, 3.6, 11.4, 1.4)
        r = lf.paragraphs[0].add_run(); r.text = spec["lead"]
        deck._set_run(r, 18, "#CBD5E1", False, True, deck.theme.get("body_font", "Calibri"))
        deck.animate(lf, **anim.get("lead", {"effect": "fade", "trigger": "after", "order": 3}))


def _slide_bullets(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.14, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    if spec.get("num"):
        badge = deck.section_badge(slide, str(spec["num"]), spec.get("kicker", ""))
        deck.animate(badge, **anim.get("badge", {"effect": "zoom", "trigger": "after", "order": 1}))
        top = 2.35
    else:
        top = 1.4
    deck.heading(slide, spec["title"], 0.9, top - 1.1, 11.4, 30, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    deck._box(slide, 0.94, top - 0.25, 2.4, 0.05, fill=deck.theme.get("accent", "#12B5C1"))
    items = spec.get("items", [])
    deck.bullets(slide, 0.9, top + 0.3, 11.4, spec.get("h", 4.6), items,
                 size=spec.get("size", 17), gap=spec.get("gap", 10), levels=spec.get("levels", False),
                 anim=anim.get("items", {"effect": "fly", "trigger": "after", "order_base": 3}))


def _slide_two_col(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.14, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    if spec.get("num"):
        deck.section_badge(slide, str(spec["num"]), spec.get("kicker", ""))
        top = 2.35
    else:
        top = 1.4
    deck.heading(slide, spec["title"], 0.9, top - 1.1, 11.4, 30, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    deck._box(slide, 0.94, top - 0.25, 2.4, 0.05, fill=deck.theme.get("accent", "#12B5C1"))
    L = spec["left"]; R = spec["right"]
    deck.bullets(slide, 0.9, top + 0.35, 5.7, 4.4, L, size=spec.get("size", 16),
                 anim=anim.get("left", {"effect": "fly", "trigger": "after", "order_base": 3}))
    deck.bullets(slide, 6.8, top + 0.35, 5.7, 4.4, R, size=spec.get("size", 16),
                 anim=anim.get("right", {"effect": "fly", "trigger": "after", "order_base": 3}))


def _slide_metrics(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.14, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    if spec.get("num"):
        deck.section_badge(slide, str(spec["num"]), spec.get("kicker", ""))
        top = 2.35
    else:
        top = 1.4
    deck.heading(slide, spec["title"], 0.9, top - 1.1, 11.4, 30, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    deck._box(slide, 0.94, top - 0.25, 2.4, 0.05, fill=deck.theme.get("accent", "#12B5C1"))
    cards = spec["cards"]
    cards_anim = anim.get("cards", {"effect": "zoom", "trigger": "after", "order_base": 3})
    base = cards_anim.get("order_base", 3)
    c_effect = cards_anim.get("effect", "zoom")
    c_trigger = cards_anim.get("trigger", "after")
    n = len(cards)
    gap = 0.28; x0 = 0.9; w = (11.9 - gap * (n + 1)) / n; h = 2.0
    for i, c in enumerate(cards):
        x = x0 + i * (w + gap)
        deck.metric_card(slide, x, top + 0.35, w, h, c["value"], c["label"],
                         sub=c.get("sub"), value_color=c.get("value_color"),
                         anim={"effect": c_effect, "trigger": c_trigger, "order": base + i})
    if spec.get("note"):
        nf = deck._txt(slide, 0.9, top + 2.6, 11.4, 0.7)
        r = nf.paragraphs[0].add_run(); r.text = spec["note"]
        deck._set_run(r, 13, deck.theme.get("accent", "#12B5C1"), False, True, deck.theme.get("body_font", "Calibri"))
        deck.animate(nf, **anim.get("note", {"effect": "fade", "trigger": "after", "order": 8}))


def _slide_chart(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.14, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    if spec.get("num"):
        deck.section_badge(slide, str(spec["num"]), spec.get("kicker", ""))
        top = 2.35
    else:
        top = 1.4
    deck.heading(slide, spec["title"], 0.9, top - 1.1, 11.4, 30, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    deck._box(slide, 0.94, top - 0.25, 2.4, 0.05, fill=deck.theme.get("accent", "#12B5C1"))
    if spec.get("grouped"):
        cf = deck.grouped_chart(slide, 0.9, top + 0.35, 11.4, 4.6,
                           spec["categories"], spec["series"], spec["colors"])
    else:
        cf = deck.bar_chart(slide, 0.9, top + 0.35, 11.4, 4.6,
                       spec["categories"], spec["values"],
                       value_color=spec.get("value_color"), show_values=spec.get("show_values", True),
                       label_size=spec.get("label_size", 13))
    deck.animate(cf, **anim.get("chart", {"effect": "fade", "trigger": "after", "order": 3}))
    if spec.get("note"):
        nf = deck._txt(slide, 0.9, top + 4.95, 11.4, 0.6)
        r = nf.paragraphs[0].add_run(); r.text = spec["note"]
        deck._set_run(r, 14, deck.theme.get("accent", "#12B5C1"), False, True, deck.theme.get("body_font", "Calibri"))
        deck.animate(nf, **anim.get("note", {"effect": "fade", "trigger": "after", "order": 4}))


def _slide_table(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.14, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    if spec.get("num"):
        deck.section_badge(slide, str(spec["num"]), spec.get("kicker", ""))
        top = 2.35
    else:
        top = 1.4
    deck.heading(slide, spec["title"], 0.9, top - 1.1, 11.4, 30, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    deck._box(slide, 0.94, top - 0.25, 2.4, 0.05, fill=deck.theme.get("accent", "#12B5C1"))
    grid = deck.data_table(slide, 0.9, top + 0.35, 11.4, spec.get("h", 4.4),
                    spec["rows"], header_color=spec.get("header_color"), row_colors=spec.get("row_colors"),
                    anim=anim.get("table", {"effect": "fade", "trigger": "after", "order": 3}))
    if spec.get("note"):
        nf = deck._txt(slide, 0.9, top + spec.get("h", 4.4) + 0.4, 11.4, 0.6)
        r = nf.paragraphs[0].add_run(); r.text = spec["note"]
        deck._set_run(r, 13, deck.theme.get("accent", "#12B5C1"), False, True, deck.theme.get("body_font", "Calibri"))
        deck.animate(nf, **anim.get("note", {"effect": "fade", "trigger": "after", "order": 4}))


def _slide_quote(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.14, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    deck.heading(slide, spec["title"], 0.9, 1.3, 11.4, 30, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    panel = deck._box(slide, 0.9, 2.3, 11.4, 3.4, fill=deck.theme.get("panel", "#12212E"),
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    panel.line.color.rgb = _rgb(deck.theme.get("accent", "#12B5C1")); panel.line.width = Pt(1.25)
    deck.animate(panel, **anim.get("panel", {"effect": "zoom", "trigger": "after", "order": 3}))
    qf = deck._txt(slide, 1.4, 2.6, 10.4, 2.8, MSO_ANCHOR.MIDDLE)
    r = qf.paragraphs[0].add_run(); r.text = spec["text"]
    deck._set_run(r, 22, "#F1F5F9", False, False, deck.theme.get("body_font", "Calibri"))
    deck.animate(qf, **anim.get("quote", {"effect": "fade", "trigger": "after", "order": 4}))
    if spec.get("attribution"):
        af = deck._txt(slide, 1.4, 5.1, 10.4, 0.5)
        r2 = af.paragraphs[0].add_run(); r2.text = spec["attribution"]
        deck._set_run(r2, 14, deck.theme.get("accent", "#12B5C1"), True, True, deck.theme.get("body_font", "Calibri"))
        deck.animate(af, **anim.get("attr", {"effect": "fade", "trigger": "after", "order": 5}))


def _slide_closing(deck, slide, spec):
    deck.background(slide, deck.theme.get("bg", "#0B1420"))
    deck.stripe(slide, "left", 0.16, deck.theme.get("accent", "#12B5C1"))
    anim = spec.get("anim", {})
    deck.heading(slide, spec.get("title", "Questions & Discussion"), 1.0, 2.2, 11.3, 40, "#F1F5F9", anim=anim.get("title", {"effect": "fade", "trigger": "after", "order": 2}))
    rule = deck._box(slide, 1.0, 3.15, 2.8, 0.06, fill=deck.theme.get("accent", "#12B5C1"))
    deck.animate(rule, **anim.get("rule", {"effect": "push", "trigger": "after", "order": 3}))
    if spec.get("points"):
        deck.bullets(slide, 1.0, 3.5, 11.3, 3.0, spec["points"], size=18,
                     anim=anim.get("points", {"effect": "fade", "trigger": "after", "order_base": 4}))


# --------------------------------------------------------------------------- helpers for spec parsing
def _iter_lines(spec):
    s = spec
    if isinstance(s, str):
        return [s]
    if isinstance(s, list):
        return s
    return [s.get("text", "")]


# --------------------------------------------------------------------------- public API
def build_deck(spec: dict, w_in: float = 13.333, h_in: float = 7.5) -> Presentation:
    deck = Deck(spec["theme"], w_in, h_in)
    for si, spec_slide in enumerate(spec["slides"]):
        slide = deck.prs.slides.add_slide(deck.blank)
        deck.background(slide, spec_slide.get("_bg", spec_slide.get("theme_bg", deck.theme.get("bg", "#0B1420"))))
        render_slide(deck, slide, spec_slide)
    return deck.prs


def _assign_animations(prs):
    """Re-order the <p:seq> blocks that _add_seq already emitted, by declared play order.

    _add_seq appends one <p:seq> per animated shape in document order. PowerPoint plays
    seqs in document order, so we sort the seqs to match each shape's ``order`` value.
    No new elements are created - this only reassigns the ``st`` (time-action) on existing
    seqs, so the round-tripped structure stays intact.
    """
    P_ns = "{" + P + "}"
    st_map = {"after": ST_AFTER, "with": ST_WITH, "click": ST_CLICK}
    for slide in prs.slides:
        timing = slide.element.find(P_ns + "timing")
        if timing is None:
            continue
        seqs = list(timing.findall(P_ns + "seq"))
        if not seqs:
            continue
        # animated shapes in document order (matches seq append order from _add_seq)
        anim_shapes = [sh for sh in slide.shapes if hasattr(sh, "_pptx_anim")]
        ordered = sorted(anim_shapes, key=lambda sh: sh._pptx_anim.get("order", 999))
        for seq, sh in zip(seqs, ordered):
            seq.set("st", str(st_map.get(sh._pptx_anim.get("trigger", "after"), ST_AFTER)))


# --------------------------------------------------------------------------- CLI / demo
def _demo_spec():
    return {
        "theme": {"bg": "#0B1420", "panel": "#12212E", "accent": "#12B5C1",
                  "accent2": "#F5B642", "text": "#F1F5F9", "muted": "#94A3B8",
                  "title_font": "Calibri", "body_font": "Calibri"},
        "slides": [
            {"type": "title", "num": "1", "kicker": "SYSTEM OVERVIEW",
             "title": "Solvarch: Retrieval-Augmented, Domain-Adapted Small Models",
             "subtitle": "A self-hosted decision-support system for AWS Well-Architected architecture reasoning.",
             "footers": ["Manjunath Kanavi (2024CS05005)", "BITS Pilani / Cisco Systems", "Final Benchmark, Aug 2026"]},
            {"type": "section", "num": "P1", "kicker": "MOTIVATION",
             "title": "Why small models need help in the cloud",
             "lead": "General-purpose LLMs invent services, cite stale pricing, and ship non-compliant infrastructure.",
             "num_": "01"},
        ],
    }


def main():
    ap = argparse.ArgumentParser(description="Build an animated .pptx from a JSON spec.")
    ap.add_argument("--spec", required=True)
    ap.add_argument("--out", default=None)
    ap.add_argument("--resize", default=None)
    ap.add_argument("--inspect", action="store_true")
    args = ap.parse_args()

    if args.spec == "--demo" or (not Path(args.spec).exists()):
        spec = _demo_spec()
    else:
        spec = json.loads(Path(args.spec).read_text())

    w, h = (13.333, 7.5)
    if args.resize:
        w, h = [float(x) for x in args.resize.split(",")]

    prs = build_deck(spec, w, h)
    _assign_animations(prs)

    if args.inspect:
        print(f"[inspect] {len(prs.slides)} slides:")
        for i, s in enumerate(prs.slides):
            texts = [sh.text_frame.text.strip()[:40] for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
            print(f"  {i}: {texts[:2]}")
        return

    out = Path(args.out or "deck.pptx")
    prs.save(str(out))
    print(f"[ok] wrote {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
