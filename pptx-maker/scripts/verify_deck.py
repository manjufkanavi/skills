#!/usr/bin/env python3
"""Verify a generated .pptx: integrity, slide count, animations, charts, text.

Usage:  python verify_deck.py <deck.pptx>

Reports per-slide shape/animation/chart counts, the first animation of each slide, and a
summary of slides with zero animations (often intentional, e.g. a cover). Requires
python-pptx (any build).
"""
import sys
import zipfile
from pathlib import Path

from pptx import Presentation

P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def verify(path: str) -> int:
    if not Path(path).exists():
        print(f"[error] not found: {path}"); return 1

    # 1. package integrity (fail fast on corruption)
    with zipfile.ZipFile(path) as z:
        bad = z.testzip()
        if bad is not None:
            print(f"[error] corrupt entry: {bad}"); return 1

    prs = Presentation(path)
    w, h = prs.slide_width / 914400, prs.slide_height / 914400
    print(f"slides: {len(prs.slides)} | size: {w:.2f}x{h:.2f}")

    total_anim = 0
    for i, s in enumerate(prs.slides):
        seqs = [e for e in s.element.iter() if e.tag.endswith("}seq")]
        charts = [sh for sh in s.shapes if sh.has_chart]
        total_anim += len(seqs)
        titles = [sh.text_frame.text.strip().split(chr(10))[0].strip()
                  for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        title = titles[0] if titles else "(no text)"
        print(f"slide {i + 1:>2}: {len(s.shapes):>2} shapes, "
              f"{len(seqs):>2} anims, {len(charts)} charts | {title[:48]}")
        for q in seqs[:4]:
            print(f"      seq type={q.get('type')} st={q.get('st')}")

    zero = sum(1 for s in prs.slides
               if not [e for e in s.element.iter() if e.tag.endswith("}seq")])
    print(f"summary: {len(prs.slides)} slides, {total_anim} total animations, "
          f"{zero} with none")
    return 0


if __name__ == "__main__":
    sys.exit(verify(sys.argv[1] if len(sys.argv) == 2 else "deck.pptx"))
